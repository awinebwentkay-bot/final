"""PPT 生成节点：生成两份演示文稿 — 主持人手卡 + 活动现场展示PPT"""

import json
import re
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from models import ActivityState

# ── 配色 ──────────────────────────────────────────────────────
C_BLUE = RGBColor(0x1A, 0x56, 0xDB)
C_DARK = RGBColor(0x1E, 0x29, 0x3B)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
C_ACCENT = RGBColor(0xE8, 0x3E, 0x8C)
C_GRAY = RGBColor(0x94, 0xA3, 0xB8)
C_GREEN = RGBColor(0x10, 0xB9, 0x81)
C_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
C_GOLD = RGBColor(0xF5, 0x9E, 0x0B)
C_TEAL = RGBColor(0x14, 0xB8, 0xA6)
C_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)

PPT_DIR = Path("PPT演示文稿")


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _textbox(slide, l, t, w, h, text, size=20, bold=False, color=C_DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    p.alignment = align
    return tb


def _ml(slide, l, t, w, h, lines, size=18, color=C_DARK, spacing=0.35):
    """多行文本框。"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(spacing * 6)
    return tb


def _header(slide, title, subtitle=""):
    """统一标题栏。"""
    _rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.9), C_BLUE)
    _textbox(slide, Inches(0.5), Inches(0.08), Inches(9), Inches(0.7),
             title, size=28, bold=True, color=C_WHITE)
    if subtitle:
        _textbox(slide, Inches(0.5), Inches(1.0), Inches(9), Inches(0.35),
                 subtitle, size=14, color=C_GRAY)


def _parse_segments(schedule: str) -> list[dict]:
    """从日程文本提取环节列表（只提取活动当天安排，跳过前期准备）。"""
    # 只取"当天安排"之后的内容
    if "当天安排" in schedule:
        schedule = schedule.split("当天安排", 1)[1]
    elif "前期准备" in schedule:
        schedule = schedule.split("前期准备", 1)[1]
        if "当天安排" in schedule:
            schedule = schedule.split("当天安排", 1)[1]

    segments = []
    for line in schedule.split("\n"):
        line = line.strip()
        if not line:
            continue
        m = re.match(r"^[-\d:：\s]*(\d{1,2}[:：]\d{2})[-~至到]?\s*(\d{1,2}[:：]\d{2})?", line)
        if m:
            time_str = m.group(0).strip("- ").strip()[:12]
            rest = line[m.end():].strip("- ").strip()
            # 去掉"（需X人）"标记
            rest = re.sub(r"（需\d+人）", "", rest).strip()
            rest = re.sub(r"\(需\d+人\)", "", rest).strip()
            if "前期准备" not in rest and "活动前" not in time_str:
                segments.append({"time": time_str, "activity": rest or line, "detail": ""})
        else:
            if "前期准备" in line or "活动前" in line:
                continue
            if segments:
                segments[-1]["detail"] += " " + line
            else:
                segments.append({"time": "", "activity": line, "detail": ""})
    return segments


def _find_title(plan: str) -> str:
    for line in plan.split("\n"):
        s = line.strip("# ").strip()
        if "主题" in s and len(s) < 30:
            t = s.split("主题")[-1].strip(""""""' ')
            if t:
                return t
        if 3 < len(s) < 25:
            return s
    return "校园活动"


def _get_confirmed_info(state: dict) -> dict:
    confirmed = state.get("poster_info_confirmed", "{}")
    try:
        return json.loads(confirmed)
    except (json.JSONDecodeError, TypeError):
        return {}


# ══════════════════════════════════════════════════════════════
# PPT-A：主持人手卡（完整主持词，活动当天逐页使用）
# ══════════════════════════════════════════════════════════════
def _build_host_card(prs: Presentation, state: dict) -> str:
    plan = state.get("activity_plan") or ""
    schedule = state.get("schedule") or ""
    host_script = state.get("host_script") or ""
    info = _get_confirmed_info(state)

    title_text = info.get("title") or _find_title(plan)
    segments = _parse_segments(schedule) or [{"time": "", "activity": "（日程待生成）", "detail": ""}]

    # 将主持稿按环节分段
    script_lines = [l.strip() for l in host_script.split("\n") if l.strip()]
    # 尝试按环节标记分段
    script_segments = []
    current_seg = []
    seg_keywords = [seg["activity"][:6] for seg in segments]

    for line in script_lines:
        # 如果该行匹配某个环节，开始新段
        matched = False
        for kw in seg_keywords:
            if kw and kw in line[:10]:
                if current_seg:
                    script_segments.append("\n".join(current_seg))
                current_seg = [line]
                matched = True
                break
        if not matched:
            current_seg.append(line)
    if current_seg:
        script_segments.append("\n".join(current_seg))

    # 如果分段数与环节数不匹配，就按环节数均分
    if len(script_segments) != len(segments):
        if len(script_lines) >= len(segments):
            chunk = max(1, len(script_lines) // len(segments))
            script_segments = []
            for i in range(len(segments)):
                start = i * chunk
                end = start + chunk if i < len(segments) - 1 else len(script_lines)
                script_segments.append("\n".join(script_lines[start:end]))
        else:
            script_segments = [host_script] * len(segments)

    # ── Slide 1: 封面 ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, C_BLUE)
    _rect(s, Inches(0), Inches(3.0), Inches(10), Inches(0.04), C_ACCENT)
    _textbox(s, Inches(0.5), Inches(1.5), Inches(9), Inches(1.2),
             title_text, size=38, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _textbox(s, Inches(0.5), Inches(3.4), Inches(9), Inches(0.5),
             f"主持人手卡  |  {info.get('date', '')} {info.get('time', '')}",
             size=18, color=C_WHITE, align=PP_ALIGN.CENTER)
    if info.get("venue"):
        _textbox(s, Inches(0.5), Inches(4.0), Inches(9), Inches(0.5),
                 f"📍 {info['venue']}", size=16, color=C_GRAY, align=PP_ALIGN.CENTER)

    # ── Slide 2: 流程总览 ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "📋 活动流程总览")
    y = Inches(1.3)
    for i, seg in enumerate(segments[:8]):
        c = C_BLUE if i % 2 == 0 else C_TEAL
        _rect(s, Inches(0.4), y, Inches(1.5), Inches(0.42), c)
        _textbox(s, Inches(0.4), y + Inches(0.03), Inches(1.5), Inches(0.36),
                 seg["time"] or f"环节{i+1}", size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _textbox(s, Inches(2.1), y + Inches(0.03), Inches(7.4), Inches(0.36),
                 seg["activity"][:35], size=16, color=C_DARK)
        y += Inches(0.55)

    # ── Slide 3+: 每环节完整主持词 ──
    for i, seg in enumerate(segments[:8]):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _header(s, f"环节 {i+1}：{seg['activity'][:20]}", seg["time"])

        # 获取该环节对应的主持词
        script_text = script_segments[i] if i < len(script_segments) else ""
        if not script_text:
            script_text = "（主持人可根据现场情况灵活发挥）"

        # 分行显示，保持可读性
        script_lines_display = script_text.split("\n")
        # 过长行做截断
        display_lines = []
        for l in script_lines_display:
            if len(l) > 60:
                # 按句号或逗号折行
                parts = re.split(r'(?<=[。，；！？])', l)
                for p in parts:
                    if p.strip():
                        display_lines.append(p.strip())
            else:
                display_lines.append(l)

        # 每页不超过 20 行，超出则分页
        page_size = 20
        for page_start in range(0, len(display_lines), page_size):
            page_lines = display_lines[page_start:page_start + page_size]
            if page_start > 0:
                s = prs.slides.add_slide(prs.slide_layouts[6])
                _header(s, f"环节 {i+1}：{seg['activity'][:20]}（续）", seg["time"])
            _ml(s, Inches(0.5), Inches(1.3), Inches(9), Inches(5.8),
                page_lines, size=16, color=C_DARK, spacing=0.3)

    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"主持人手卡_{ts}.pptx"
    path = PPT_DIR / filename
    prs.save(str(path))
    return str(path)


# ══════════════════════════════════════════════════════════════
# PPT-B：活动现场展示PPT（大屏投影用）
# ══════════════════════════════════════════════════════════════
def _build_display_ppt(prs: Presentation, state: dict) -> str:
    plan = state.get("activity_plan", "")
    schedule = state.get("schedule", "")
    info = _get_confirmed_info(state)

    title_text = info.get("title") or _find_title(plan)
    segments = _parse_segments(schedule) or [{"time": "", "activity": "（日程待生成）", "detail": ""}]

    # ── Slide 1: 封面 ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, C_DARK)
    # 装饰线条
    _rect(s, Inches(0), Inches(3.0), Inches(10), Inches(0.04), C_BLUE)
    _rect(s, Inches(0), Inches(3.1), Inches(10), Inches(0.02), C_ACCENT)
    # 大字标题
    _textbox(s, Inches(1), Inches(1.5), Inches(8), Inches(1.5),
             title_text, size=46, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # 信息
    info_lines = []
    if info.get("date"):
        info_lines.append(f"📅 {info['date']}")
    if info.get("time"):
        info_lines.append(f"⏰ {info['time']}")
    if info.get("venue"):
        info_lines.append(f"📍 {info['venue']}")
    if info_lines:
        _textbox(s, Inches(1), Inches(3.5), Inches(8), Inches(0.8),
                 "  |  ".join(info_lines), size=20, color=C_GRAY, align=PP_ALIGN.CENTER)
    # 主办方
    if info.get("organizer") and info["organizer"] != "待定":
        _textbox(s, Inches(1), Inches(4.5), Inches(8), Inches(0.5),
                 f"主办：{info['organizer']}", size=16, color=C_GRAY, align=PP_ALIGN.CENTER)

    # ── Slide 2: 欢迎页 ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, C_BLUE)
    _textbox(s, Inches(1), Inches(2.0), Inches(8), Inches(1.2),
             "欢迎参加", size=26, color=C_WHITE, align=PP_ALIGN.CENTER)
    _textbox(s, Inches(1), Inches(3.0), Inches(8), Inches(1.5),
             title_text, size=42, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _textbox(s, Inches(1), Inches(4.8), Inches(8), Inches(0.6),
             "请将手机调至静音，享受本次活动", size=18, color=C_GRAY, align=PP_ALIGN.CENTER)

    # ── Slide 3: 流程总览 ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, C_LIGHT)
    _rect(s, Inches(0), Inches(0), Inches(10), Inches(0.9), C_BLUE)
    _textbox(s, Inches(0.5), Inches(0.1), Inches(9), Inches(0.7),
             "活动流程", size=30, bold=True, color=C_WHITE)

    y = Inches(1.3)
    for i, seg in enumerate(segments[:7]):
        c = C_BLUE if i % 2 == 0 else C_TEAL
        _rect(s, Inches(0.6), y, Inches(1.8), Inches(0.5), c)
        _textbox(s, Inches(0.6), y + Inches(0.05), Inches(1.8), Inches(0.4),
                 seg["time"] or f"0{i+1}", size=17, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _textbox(s, Inches(2.6), y + Inches(0.05), Inches(7), Inches(0.4),
                 seg["activity"][:40], size=18, bold=True, color=C_DARK)
        y += Inches(0.65)

    # ── Slide 4+: 每环节过渡页（大屏展示） ──
    for i, seg in enumerate(segments[:6]):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        # 左半色块
        _rect(s, Inches(0), Inches(0), Inches(4), Inches(7.5), C_BLUE)
        # 环节编号
        _textbox(s, Inches(0.5), Inches(1.5), Inches(3), Inches(0.8),
                 f"0{i+1}", size=50, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # 环节名称
        _textbox(s, Inches(0.5), Inches(2.5), Inches(3), Inches(0.6),
                 seg["activity"][:15], size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # 时间
        if seg["time"]:
            _textbox(s, Inches(0.5), Inches(3.3), Inches(3), Inches(0.5),
                     seg["time"], size=16, color=C_GRAY, align=PP_ALIGN.CENTER)
        # 右侧大字
        _textbox(s, Inches(4.5), Inches(2.5), Inches(5), Inches(2),
                 seg["activity"][:20], size=38, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

    # ── 结尾页 ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, C_DARK)
    _rect(s, Inches(0), Inches(3.5), Inches(10), Inches(0.04), C_BLUE)
    _textbox(s, Inches(1), Inches(2.0), Inches(8), Inches(1.5),
             "感谢参与", size=46, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if info.get("organizer") and info["organizer"] != "待定":
        _textbox(s, Inches(1), Inches(4.0), Inches(8), Inches(0.6),
                 f"主办：{info['organizer']}", size=18, color=C_GRAY, align=PP_ALIGN.CENTER)

    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"活动现场展示_{ts}.pptx"
    path = PPT_DIR / filename
    prs.save(str(path))
    return str(path)


# ══════════════════════════════════════════════════════════════
# 入口
# ══════════════════════════════════════════════════════════════
def ppt_agent(state: ActivityState) -> ActivityState:
    """生成两份 PPT：主持人手卡 + 活动现场展示PPT。"""
    PPT_DIR.mkdir(exist_ok=True)

    # ── PPT-A：主持人手卡 ──
    print(f"[PPT] 正在生成主持人手卡...", flush=True)
    prs_a = Presentation()
    prs_a.slide_width = Inches(10)
    prs_a.slide_height = Inches(7.5)
    path_a = _build_host_card(prs_a, state)
    print(f"[PPT] 主持人手卡已生成：{path_a}", flush=True)

    # ── PPT-B：活动现场展示PPT ──
    print(f"[PPT] 正在生成活动现场展示PPT...", flush=True)
    prs_b = Presentation()
    prs_b.slide_width = Inches(10)
    prs_b.slide_height = Inches(7.5)
    path_b = _build_display_ppt(prs_b, state)

    # 合并路径信息
    state["ppt_path"] = f"{path_a}\n    🖥️  {path_b}"
    state["log"].append(f"【PPT】主持人手卡 + 活动现场展示PPT 已生成")
    print(f"[PPT] 活动现场展示PPT已生成：{path_b}", flush=True)

    return state